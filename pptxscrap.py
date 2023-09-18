from pptx import Presentation

# Function to extract text from a PowerPoint (.pptx) file
def extract_text_from_pptx(pptx_path):
    try:
        presentation = Presentation(pptx_path)
        text_content = []

        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)

        return "\n".join(text_content)

    except Exception as e:
        print("Error extracting text from PowerPoint:", e)
        return None

# Example usage
if __name__ == "__main__":
    pptx_path = "C:/Users/Admin/Desktop/PIB_Report.pptx"  # Replace with the path to your PowerPoint file
    text_data = extract_text_from_pptx(pptx_path)

    if text_data:
        print(text_data)
